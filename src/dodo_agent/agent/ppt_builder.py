import asyncio
import json
import logging
import os
import re
import uuid
from datetime import datetime
from io import BytesIO
from pathlib import Path

import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu
from tavily import TavilyClient

from src.dodo_agent.agent.react_agent import build_llm
from src.dodo_agent.config.settings import settings
from src.dodo_agent.storage.db import new_session
from src.dodo_agent.storage.models.ai_ppt_inst import AiPptInst, PptInstRepo

logger = logging.getLogger("dodo")

STATE_ORDER = ["INIT", "SCHEMA", "OUTLINE", "CONTENT", "RENDER", "SUCCESS"]

JSON_PATTERN = re.compile(r"```(?:json)?\s*\n?(.*?)\n?```", re.DOTALL)


class _AgentStopped(Exception):
    pass


# ---- helpers ----

def _make_event(event_type: str, **kwargs) -> dict:
    payload = {"type": event_type}
    if event_type == "thinking" and "content" in kwargs:
        kwargs["content"] = kwargs["content"] + "\n\n"
    payload.update(kwargs)
    return {"event": "message", "data": json.dumps(payload, ensure_ascii=False)}


def _make_sse(text: str) -> dict:
    return {"event": "message", "data": text}


def _parse_json(text: str) -> dict | None:
    m = JSON_PATTERN.search(text)
    json_str = m.group(1) if m else text
    try:
        return json.loads(json_str)
    except json.JSONDecodeError:
        pass
    match = re.search(r"\{.*\}", text, re.DOTALL)
    if match:
        try:
            return json.loads(match.group(0))
        except json.JSONDecodeError:
            pass
    return None


def _load_or_create_inst(conversation_id: str, query: str) -> AiPptInst:
    repo = PptInstRepo()
    existing = repo.find_by_conversation_id(conversation_id)
    if existing and existing.status not in ("SUCCESS", "FAILED"):
        logger.info(f"恢复PPT任务: {conversation_id}, 状态={existing.status}")
        return existing
    if existing and existing.status in ("SUCCESS", "FAILED"):
        repo.delete(existing)
    inst = AiPptInst(
        conversation_id=conversation_id,
        query=query,
        status="INIT",
        create_time=datetime.now(),
        update_time=datetime.now(),
    )
    repo.save(inst)
    return inst


def _save_inst(inst: AiPptInst):
    inst.update_time = datetime.now()
    repo = PptInstRepo()
    repo._s.merge(inst)
    repo._s.flush()
    if repo._own_session:
        repo._s.commit()
    repo.close()


# ---- prompts ----

INIT_PROMPT = """你是一个专业的PPT设计顾问。分析用户的需求，生成一份结构清晰的PPT需求说明。

用户需求：{query}

请分析并输出JSON格式：
```json
{{
  "intent": "CREATE_PPT",
  "title": "PPT主标题（简洁有力，不超过15字）",
  "requirement": "需求澄清描述（50字以内）",
  "pageCount": 8,
  "style": "科技/商务/简约/学术/创意 中选一个"
}}
```

规则：
1. pageCount 建议8-15页
2. style 根据主题智能判断
3. 只输出JSON，不要其他文字"""

SCHEMA_PROMPT = """你是一个PPT结构设计师。根据需求规划每页的页面类型和标题。

PPT标题：{title}
需求：{requirement}
预计页数：约{page_count}页
风格：{style}

输出JSON格式，slides数组中每页包含：index(从1开始)、pageType(COVER/CATALOG/CONTENT/COMPARE/END)、title、description

```json
{{
  "slides": [
    {{"index": 1, "pageType": "COVER", "title": "封面标题", "description": "副标题/作者信息"}},
    {{"index": 2, "pageType": "CATALOG", "title": "目录", "description": "列出3-5个章节"}},
    {{"index": 3, "pageType": "CONTENT", "title": "内容页标题", "description": "这页要讲什么"}}
  ]
}}
```

pageType说明：
- COVER: 封面页（第1页）
- CATALOG: 目录页（第2页）
- CONTENT: 正文内容页
- COMPARE: 对比页（左右对比）
- END: 结束页（最后1页）

规则：第1页必须是COVER，第2页建议CATALOG，最后1页必须是END，中间根据内容安排CONTENT/COMPARE。
只输出JSON，不要其他文字。"""

OUTLINE_PROMPT = """你是一个PPT大纲撰写专家。为每页生成详细的大纲。

PPT主题：{title}
页面结构：{schema_json}

为每页生成详细大纲，输出JSON：
```json
{{
  "slides": [
    {{
      "index": 1,
      "title": "原标题",
      "keyPoints": ["要点1", "要点2", "要点3"],
      "imageKeywords": "用于搜索配图的英文关键词"
    }}
  ]
}}
```

规则：
1. 每页3-5个要点
2. imageKeywords 用英文，适合图片搜索引擎
3. COVER和END页imageKeywords可为空字符串
4. 只输出JSON"""

CONTENT_PROMPT = """你是一个PPT内容撰写专家。为指定页面生成完整的演示内容。

PPT主题：{theme}
页面类型：{page_type}
页面标题：{title}
大纲要点：{key_points}

请生成该页的完整内容，输出JSON：
```json
{{
  "title": "页面标题",
  "subtitle": "副标题（可选，CONTENT页可为空字符串）",
  "bullets": ["要点1的详细内容", "要点2的详细内容"],
  "notes": "演讲备注（可选）"
}}
```

对于COMPARE类型页面，额外包含leftTitle/leftBullets/rightTitle/rightBullets。
规则：bullets中每条20-40字，简洁有力。只输出JSON。"""


# ---- Tavily image search ----

async def _search_images(keywords: str, max_count: int = 2) -> list[str]:
    if not keywords or not settings.tavily_api_key:
        return []
    try:
        client = TavilyClient(api_key=settings.tavily_api_key)
        response = await asyncio.to_thread(
            client.search, keywords, include_images=True, max_results=max_count
        )
        images = response.get("images", [])
        return images[:max_count] if images else []
    except Exception as e:
        logger.warning(f"Tavily图片搜索失败 ({keywords}): {e}")
        return []


async def _download_image(url: str) -> bytes | None:
    try:
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(url)
            if resp.status_code == 200 and len(resp.content) > 0:
                return resp.content
    except Exception as e:
        logger.warning(f"图片下载失败 {url[:80]}: {e}")
    return None


# ---- PPTX renderer ----

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

PRIMARY_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_COLOR = RGBColor(0x00, 0x72, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
ACCENT_ORANGE = RGBColor(0xE9, 0x45, 0x60)
GRAY_TEXT = RGBColor(0x88, 0x88, 0x88)

FONT_TITLE = "Arial"
FONT_BODY = "Arial"


def _add_blank_slide(prs: Presentation) -> object:
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def _set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text: str,
                 font_size: int = 18, bold: bool = False, color: RGBColor = DARK_TEXT,
                 alignment: int = PP_ALIGN.LEFT, font_name: str = FONT_BODY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_bullet_frame(slide, left, top, width, height, bullets: list[str],
                      font_size: int = 16, color: RGBColor = DARK_TEXT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {bullet}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = Pt(8)
        p.level = 0
    return tf


def _add_image_safe(slide, image_bytes: bytes, left, top, width, height) -> bool:
    try:
        stream = BytesIO(image_bytes)
        slide.shapes.add_picture(stream, Inches(left), Inches(top), Inches(width), Inches(height))
        return True
    except Exception as e:
        logger.warning(f"插入图片失败: {e}")
        return False


def _add_accent_bar(slide, left, top, width, height, color: RGBColor = ACCENT_COLOR):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def render_pptx(slides_data: list[dict], title: str, output_path: str):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for s in slides_data:
        page_type = s.get("pageType", "CONTENT")

        if page_type == "COVER":
            _render_cover(prs, s, title)
        elif page_type == "CATALOG":
            _render_catalog(prs, s)
        elif page_type == "COMPARE":
            _render_compare(prs, s)
        elif page_type == "END":
            _render_end(prs, s)
        else:
            _render_content(prs, s)

    prs.save(output_path)
    logger.info(f"PPTX已生成: {output_path}, 共{len(slides_data)}页")


def _render_cover(prs, s, title):
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, PRIMARY_COLOR)

    _add_accent_bar(slide, 1.5, 3.2, 1.2, 0.06, ACCENT_COLOR)

    _add_textbox(slide, 1.5, 2.2, 10.5, 1.2,
                 s.get("title", title), font_size=44, bold=True, color=WHITE,
                 alignment=PP_ALIGN.LEFT)

    desc = s.get("description", "")
    if desc:
        _add_textbox(slide, 1.5, 3.6, 10.5, 0.8,
                     desc, font_size=20, color=RGBColor(0xBB, 0xBB, 0xCC),
                     alignment=PP_ALIGN.LEFT)

    _add_textbox(slide, 1.5, 5.8, 10.5, 0.5,
                 "AI Generated by Dodo Agent", font_size=12, color=GRAY_TEXT)


def _render_catalog(prs, s):
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)

    _add_accent_bar(slide, 1.0, 1.0, 0.08, 0.6, ACCENT_COLOR)
    _add_textbox(slide, 1.3, 0.9, 8, 0.7,
                 s.get("title", "目录"), font_size=36, bold=True, color=PRIMARY_COLOR)

    items = s.get("items", [])
    if not items and s.get("keyPoints"):
        items = s["keyPoints"]
    if not items and s.get("bullets"):
        items = s["bullets"]

    for i, item in enumerate(items):
        y = 2.2 + i * 1.0
        num = f"0{i + 1}" if i < 9 else str(i + 1)
        _add_textbox(slide, 1.5, y, 0.8, 0.6,
                     num, font_size=32, bold=True, color=ACCENT_COLOR)
        text = item if isinstance(item, str) else item.get("title", str(item))
        _add_textbox(slide, 2.5, y + 0.05, 9, 0.5,
                     text, font_size=22, color=DARK_TEXT)


def _render_content(prs, s):
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)

    _add_accent_bar(slide, 1.0, 0.0, 11.333, 0.04, ACCENT_COLOR)

    _add_textbox(slide, 1.0, 0.5, 11.3, 0.7,
                 s.get("title", ""), font_size=32, bold=True, color=PRIMARY_COLOR)

    subtitle = s.get("subtitle", "")
    if subtitle:
        _add_textbox(slide, 1.0, 1.15, 11.3, 0.5,
                     subtitle, font_size=16, color=GRAY_TEXT)

    has_image = bool(s.get("_image_path") and os.path.exists(s["_image_path"]))
    text_width = 6.5 if has_image else 11.3

    bullets = s.get("bullets", []) or s.get("keyPoints", [])
    if bullets:
        _add_bullet_frame(slide, 1.0, 1.8, text_width, 4.8,
                          bullets, font_size=18, color=DARK_TEXT)

    if has_image:
        with open(s["_image_path"], "rb") as f:
            _add_image_safe(slide, f.read(), 8.0, 1.8, 4.5, 4.5)

    notes = s.get("notes", "")
    if notes:
        _add_textbox(slide, 1.0, 6.8, 11.3, 0.4,
                     notes, font_size=11, color=GRAY_TEXT)


def _render_compare(prs, s):
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)

    _add_accent_bar(slide, 1.0, 0.0, 11.333, 0.04, ACCENT_COLOR)
    _add_textbox(slide, 1.0, 0.5, 11.3, 0.7,
                 s.get("title", ""), font_size=32, bold=True, color=PRIMARY_COLOR)

    left_title = s.get("leftTitle", "方案A")
    right_title = s.get("rightTitle", "方案B")
    left_bullets = s.get("leftBullets", []) or s.get("leftContent", [])
    right_bullets = s.get("rightBullets", []) or s.get("rightContent", [])

    _add_accent_bar(slide, 1.0, 1.6, 5.0, 0.04, ACCENT_COLOR)
    _add_textbox(slide, 1.0, 1.8, 5.0, 0.5,
                 left_title, font_size=24, bold=True, color=ACCENT_COLOR)

    if left_bullets:
        _add_bullet_frame(slide, 1.0, 2.4, 5.0, 4.2,
                          left_bullets, font_size=15)

    _add_accent_bar(slide, 7.0, 1.6, 5.0, 0.04, ACCENT_ORANGE)
    _add_textbox(slide, 7.0, 1.8, 5.0, 0.5,
                 right_title, font_size=24, bold=True, color=ACCENT_ORANGE)

    if right_bullets:
        _add_bullet_frame(slide, 7.0, 2.4, 5.0, 4.2,
                          right_bullets, font_size=15)

    _add_textbox(slide, 6.3, 3.5, 0.8, 0.6,
                 "VS", font_size=28, bold=True, color=GRAY_TEXT,
                 alignment=PP_ALIGN.CENTER)


def _render_end(prs, s):
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, PRIMARY_COLOR)

    _add_textbox(slide, 1.5, 2.8, 10.5, 1.2,
                 s.get("title", "谢谢观看"), font_size=48, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    desc = s.get("description", "")
    if desc:
        _add_textbox(slide, 1.5, 4.2, 10.5, 0.8,
                     desc, font_size=20, color=GRAY_TEXT,
                     alignment=PP_ALIGN.CENTER)


# ---- state handlers ----

async def _handle_init(inst: AiPptInst, llm, cancel_event: asyncio.Event):
    yield _make_event("thinking", content="正在分析PPT需求...")

    response = await llm.ainvoke([
        ("system", INIT_PROMPT.format(query=inst.query)),
    ])
    text = response.content if hasattr(response, "content") else str(response)
    data = _parse_json(text) or {}

    inst.requirement = data.get("requirement", inst.query)
    inst.template_code = data.get("style", "auto")
    inst.ppt_schema = json.dumps(data, ensure_ascii=False)
    inst.status = "SCHEMA"
    _save_inst(inst)

    title = data.get("title", "")
    page_count = data.get("pageCount", 8)
    style = data.get("style", "auto")
    yield _make_event("thinking",
                      content=f"需求分析完成：主题「{title}」，{page_count}页，风格：{style}")


async def _handle_schema(inst: AiPptInst, llm, cancel_event: asyncio.Event):
    yield _make_event("thinking", content="正在规划PPT结构...")

    try:
        init_data = json.loads(inst.ppt_schema or "{}")
    except json.JSONDecodeError:
        init_data = {}
    title = init_data.get("title", inst.query[:15])
    page_count = init_data.get("pageCount", 8)
    style = init_data.get("style", "auto")

    response = await llm.ainvoke([
        ("system", SCHEMA_PROMPT.format(
            title=title, requirement=inst.requirement,
            page_count=page_count, style=style)),
    ])
    text = response.content if hasattr(response, "content") else str(response)
    data = _parse_json(text) or {"slides": []}
    slides = data.get("slides", [])

    schema_data = {"title": title, "style": style, "slides": slides}
    inst.ppt_schema = json.dumps(schema_data, ensure_ascii=False)
    inst.status = "OUTLINE"
    _save_inst(inst)

    yield _make_event("thinking",
                      content=f"PPT结构规划完成，共 {len(slides)} 页")


async def _handle_outline(inst: AiPptInst, llm, cancel_event: asyncio.Event):
    yield _make_event("thinking", content="正在生成详细大纲...")

    schema_json = inst.ppt_schema or "{}"
    try:
        schema_data = json.loads(schema_json)
    except json.JSONDecodeError:
        schema_data = {"title": inst.query, "slides": []}

    response = await llm.ainvoke([
        ("system", OUTLINE_PROMPT.format(
            title=schema_data.get("title", ""),
            schema_json=schema_json)),
    ])
    text = response.content if hasattr(response, "content") else str(response)
    data = _parse_json(text) or {"slides": []}
    outline_slides = data.get("slides", [])

    inst.outline = json.dumps(outline_slides, ensure_ascii=False)
    inst.status = "CONTENT"
    _save_inst(inst)

    yield _make_event("thinking",
                      content=f"大纲生成完成，共 {len(outline_slides)} 页")


async def _handle_content(inst: AiPptInst, llm, cancel_event: asyncio.Event):
    try:
        schema_data = json.loads(inst.ppt_schema or "{}")
    except json.JSONDecodeError:
        schema_data = {"title": inst.query, "slides": []}
    try:
        outline_slides = json.loads(inst.outline or "[]")
    except json.JSONDecodeError:
        outline_slides = []

    theme = schema_data.get("title", "")
    schema_slides = schema_data.get("slides", [])
    total = len(schema_slides)

    yield _make_event("thinking", content=f"正在填充幻灯片内容，共 {total} 页...")

    full_slides: list[dict] = []
    for i, schema_slide in enumerate(schema_slides):
        if cancel_event.is_set():
            raise _AgentStopped

        slide_index = schema_slide.get("index", i + 1)
        page_type = schema_slide.get("pageType", "CONTENT")
        slide_title = schema_slide.get("title", "")

        yield _make_event("thinking",
                          content=f"[{i + 1}/{total}] 生成中：{slide_title}")

        outline_match = next(
            (o for o in outline_slides if o.get("index") == slide_index), {})
        key_points = outline_match.get("keyPoints", [])
        image_kw = outline_match.get("imageKeywords", "")

        slide_data = {
            "index": slide_index,
            "pageType": page_type,
            "title": slide_title,
            "description": schema_slide.get("description", ""),
        }

        if page_type in ("CONTENT", "COMPARE"):
            response = await llm.ainvoke([
                ("system", CONTENT_PROMPT.format(
                    theme=theme, page_type=page_type,
                    title=slide_title,
                    key_points=json.dumps(key_points, ensure_ascii=False))),
            ])
            text = response.content if hasattr(response, "content") else str(response)
            content_data = _parse_json(text) or {}

            slide_data["bullets"] = content_data.get("bullets", key_points)
            slide_data["subtitle"] = content_data.get("subtitle", "")
            slide_data["notes"] = content_data.get("notes", "")
            if page_type == "COMPARE":
                slide_data["leftTitle"] = content_data.get("leftTitle", "")
                slide_data["leftBullets"] = content_data.get("leftBullets", [])
                slide_data["rightTitle"] = content_data.get("rightTitle", "")
                slide_data["rightBullets"] = content_data.get("rightBullets", [])
        elif page_type == "CATALOG":
            items = outline_match.get("keyPoints", [])
            slide_data["items"] = [str(item) for item in items] if items else [slide_title]
        elif page_type == "COVER":
            slide_data["title"] = theme
            slide_data["description"] = schema_slide.get("description", "")
        elif page_type == "END":
            slide_data["title"] = "谢谢观看"
            slide_data["description"] = theme

        if image_kw and page_type in ("CONTENT",):
            images = await _search_images(image_kw, max_count=1)
            if images:
                img_bytes = await _download_image(images[0])
                if img_bytes:
                    img_dir = Path(settings.upload_dir) / "pptx_images"
                    img_dir.mkdir(parents=True, exist_ok=True)
                    img_name = f"{uuid.uuid4().hex[:12]}.png"
                    img_path = img_dir / img_name
                    img_path.write_bytes(img_bytes)
                    slide_data["_image_path"] = str(img_path)

        full_slides.append(slide_data)

        yield _make_event("thinking",
                          content=f"[{i + 1}/{total}] 已完成：{slide_title}")

    inst.ppt_schema = json.dumps({
        "title": theme,
        "style": schema_data.get("style", "auto"),
        "slides": schema_slides,
    }, ensure_ascii=False)
    inst.outline = json.dumps(full_slides, ensure_ascii=False)
    inst.status = "RENDER"
    _save_inst(inst)

    yield _make_event("thinking",
                      content=f"内容填充完成，共 {len(full_slides)} 页")


async def _handle_render(inst: AiPptInst, llm, cancel_event: asyncio.Event):
    yield _make_event("thinking", content="正在渲染PPT文件...")

    full_slides = json.loads(inst.outline or "[]")
    ppt_schema = json.loads(inst.ppt_schema or "{}")
    title = ppt_schema.get("title", inst.query)

    output_dir = Path(settings.upload_dir) / "pptx"
    output_dir.mkdir(parents=True, exist_ok=True)
    file_name = f"{uuid.uuid4().hex[:12]}.pptx"
    output_path = output_dir / file_name

    await asyncio.to_thread(render_pptx, full_slides, title, str(output_path))

    for s in full_slides:
        img_path = s.get("_image_path", "")
        if img_path and os.path.exists(img_path):
            try:
                os.remove(img_path)
            except OSError:
                pass

    from minio import Minio
    from minio.error import S3Error

    file_url = ""
    if settings.minio_endpoint:
        try:
            client = Minio(
                settings.minio_endpoint,
                access_key=settings.minio_access_key,
                secret_key=settings.minio_secret_key,
                secure=False,
            )
            obj_name = f"pptx/{file_name}"
            client.fput_object(settings.minio_bucket, obj_name, str(output_path))
            file_url = f"minio://{settings.minio_bucket}/{obj_name}"
        except S3Error as e:
            logger.error(f"MinIO上传PPT失败: {e}")

    if not file_url:
        file_url = f"local://{output_path}"

    inst.file_url = file_url
    inst.status = "SUCCESS"
    _save_inst(inst)

    yield _make_event("thinking", content="PPT文件渲染完成")


STATE_HANDLERS = {
    "INIT": _handle_init,
    "SCHEMA": _handle_schema,
    "OUTLINE": _handle_outline,
    "CONTENT": _handle_content,
    "RENDER": _handle_render,
}


# ---- main stream ----

async def ppt_builder_stream(conversation_id: str, query: str, cancel_event: asyncio.Event):
    llm = build_llm()

    try:
        inst = _load_or_create_inst(conversation_id, query)

        start_state = inst.status or "INIT"
        if start_state not in STATE_ORDER:
            start_state = "INIT"

        start_idx = STATE_ORDER.index(start_state)
        logger.info(f"PPT Builder 启动: conv={conversation_id}, 起始状态={start_state}")

        for i in range(start_idx, len(STATE_ORDER)):
            state = STATE_ORDER[i]
            if state == "SUCCESS":
                schema = json.loads(inst.ppt_schema or "{}")
                title = schema.get("title", "output")
                download_path = f"/agent/pptx/download?conversationId={conversation_id}"
                yield _make_event("text",
                                  content=f"\n\nPPT已生成：[点击下载 {title}.pptx]({download_path})")
                yield _make_event("file_ready",
                                  url=inst.file_url or "",
                                  fileName=f"{title}.pptx",
                                  message="PPT已生成，点击下载")
                break

            if cancel_event.is_set():
                raise _AgentStopped

            handler = STATE_HANDLERS.get(state)
            if not handler:
                inst.status = "FAILED"
                inst.error_msg = f"未知状态: {state}"
                _save_inst(inst)
                yield _make_event("error", message=inst.error_msg)
                break

            try:
                async for event in handler(inst, llm, cancel_event):
                    yield event
                inst = PptInstRepo().find_by_conversation_id(conversation_id) or inst
            except _AgentStopped:
                raise
            except Exception as e:
                logger.error(f"PPT状态 {state} 异常: {e}")
                inst.status = "FAILED"
                inst.error_msg = str(e)
                _save_inst(inst)
                yield _make_event("error", message=f"PPT生成失败({state}): {e}")
                break

    except _AgentStopped:
        yield _make_sse(json.dumps({"type": "error", "content": "用户已停止"}, ensure_ascii=False))
    except Exception as e:
        logger.error(f"PPT Builder异常: {e}")
        yield _make_sse(json.dumps({"type": "error", "content": str(e)}, ensure_ascii=False))
