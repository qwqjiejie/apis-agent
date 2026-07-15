import asyncio
import json
import os
import re
import uuid
from datetime import datetime
from io import BytesIO
from pathlib import Path

import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu
from tavily import TavilyClient

from src.dodo_agent.agent.base_agent import BaseAgent
from src.dodo_agent.common.llm import build_llm
from src.dodo_agent.common.logger import logger
from src.dodo_agent.common.streaming import AgentStopped, make_event, make_sse
from src.dodo_agent.config.settings import get_settings
from src.dodo_agent.storage.models.ai_ppt_inst import AiPptInst, PptInstRepo

# =============================================================================
# 常量
# =============================================================================

# 状态机顺序：INIT → SCHEMA → OUTLINE → CONTENT → RENDER → SUCCESS
STATE_ORDER = ["INIT", "SCHEMA", "OUTLINE", "CONTENT", "RENDER", "SUCCESS"]

JSON_PATTERN = re.compile(r"```(?:json)?\s*\n?(.*?)\n?```", re.DOTALL)


# =============================================================================
# 工具函数
# =============================================================================

def _parse_json(text: str) -> dict | None:
    """从 LLM 输出中解析 JSON。两层容错：
    1. 优先匹配 ```json ... ``` 代码块
    2. 回退到匹配最外层 { ... } 对象
    """
    m = JSON_PATTERN.search(text)
    json_str = m.group(1) if m else text
    try:
        return json.loads(json_str)
    except json.JSONDecodeError:
        pass
    match = re.search(r"\{.*\}", text, re.DOTALL)
    if match:
        try:
            return json.loads(match.group(0))
        except json.JSONDecodeError:
            pass
    return None


# =============================================================================
# LLM Prompts — 五阶段提示词
# =============================================================================

INIT_PROMPT = """你是一个专业的PPT设计顾问。分析用户的需求，生成一份结构清晰的PPT需求说明。

用户需求：{query}

请分析并输出JSON格式：
```json
{{
  "intent": "CREATE_PPT",
  "title": "PPT主标题（简洁有力，不超过15字）",
  "requirement": "需求澄清描述（50字以内）",
  "pageCount": 8,
  "style": "科技/商务/简约/学术/创意 中选一个"
}}
```

规则：
1. pageCount 建议8-15页
2. style 根据主题智能判断
3. 只输出JSON，不要其他文字"""

SCHEMA_PROMPT = """你是一个PPT结构设计师。根据需求规划每页的页面类型和标题。

PPT标题：{title}
需求：{requirement}
预计页数：约{page_count}页
风格：{style}

输出JSON格式，slides数组中每页包含：index(从1开始)、pageType(COVER/CATALOG/CONTENT/COMPARE/END)、title、description

```json
{{
  "slides": [
    {{"index": 1, "pageType": "COVER", "title": "封面标题", "description": "副标题/作者信息"}},
    {{"index": 2, "pageType": "CATALOG", "title": "目录", "description": "列出3-5个章节"}},
    {{"index": 3, "pageType": "CONTENT", "title": "内容页标题", "description": "这页要讲什么"}}
  ]
}}
```

pageType说明：
- COVER: 封面页（第1页）
- CATALOG: 目录页（第2页）
- CONTENT: 正文内容页
- COMPARE: 对比页（左右对比）
- END: 结束页（最后1页）

规则：第1页必须是COVER，第2页建议CATALOG，最后1页必须是END，中间根据内容安排CONTENT/COMPARE。
只输出JSON，不要其他文字。"""

OUTLINE_PROMPT = """你是一个PPT大纲撰写专家。为每页生成详细的大纲。

PPT主题：{title}
页面结构：{schema_json}

为每页生成详细大纲，输出JSON：
```json
{{
  "slides": [
    {{
      "index": 1,
      "title": "原标题",
      "keyPoints": ["要点1", "要点2", "要点3"],
      "imageKeywords": "用于搜索配图的英文关键词"
    }}
  ]
}}
```

规则：
1. 每页3-5个要点
2. imageKeywords 用英文，适合图片搜索引擎
3. COVER和END页imageKeywords可为空字符串
4. 只输出JSON"""

CONTENT_PROMPT = """你是一个PPT内容撰写专家。为指定页面生成完整的演示内容。

PPT主题：{theme}
页面类型：{page_type}
页面标题：{title}
大纲要点：{key_points}

请生成该页的完整内容，输出JSON：
```json
{{
  "title": "页面标题",
  "subtitle": "副标题（可选，CONTENT页可为空字符串）",
  "bullets": ["要点1的详细内容", "要点2的详细内容"],
  "notes": "演讲备注（可选）"
}}
```

对于COMPARE类型页面，额外包含leftTitle/leftBullets/rightTitle/rightBullets。
规则：bullets中每条20-40字，简洁有力。只输出JSON。"""


# =============================================================================
# Tavily 图片搜索
# =============================================================================

async def _search_images(keywords: str, max_count: int = 2) -> list[str]:
    """通过 Tavily API 搜索配图，返回图片 URL 列表。"""
    if not keywords or not get_settings().tavily_api_key:
        return []
    try:
        client = TavilyClient(api_key=get_settings().tavily_api_key)
        response = await asyncio.to_thread(
            client.search, keywords, include_images=True, max_results=max_count
        )
        images = response.get("images", [])
        return images[:max_count] if images else []
    except Exception as e:
        logger.warning(f"Tavily图片搜索失败 ({keywords}): {e}")
        return []


async def _download_image(url: str) -> bytes | None:
    """下载图片到内存，15 秒超时。"""
    try:
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(url)
            if resp.status_code == 200 and len(resp.content) > 0:
                return resp.content
    except Exception as e:
        logger.warning(f"图片下载失败 {url[:80]}: {e}")
    return None


# =============================================================================
# PPTX 渲染引擎 — python-pptx 底层操作
# =============================================================================

SLIDE_WIDTH = Inches(13.333)    # 宽屏 16:9
SLIDE_HEIGHT = Inches(7.5)

# 配色方案
PRIMARY_COLOR = RGBColor(0x1A, 0x1A, 0x2E)     # 深蓝黑（封面/结束页背景）
ACCENT_COLOR = RGBColor(0x00, 0x72, 0xD8)       # 亮蓝（强调色）
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
ACCENT_ORANGE = RGBColor(0xE9, 0x45, 0x60)      # 橙红（对比页右侧强调）
GRAY_TEXT = RGBColor(0x88, 0x88, 0x88)

FONT_TITLE = "Arial"
FONT_BODY = "Arial"


def _add_blank_slide(prs: Presentation) -> object:
    """创建空白幻灯片（layout index 6 = blank layout）。"""
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def _set_slide_bg(slide, color: RGBColor):
    """设置幻灯片纯色背景。"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text: str,
                 font_size: int = 18, bold: bool = False, color: RGBColor = DARK_TEXT,
                 alignment: int = PP_ALIGN.LEFT, font_name: str = FONT_BODY):
    """在幻灯片上添加文本框，返回 text_frame 供后续操作。"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_bullet_frame(slide, left, top, width, height, bullets: list[str],
                      font_size: int = 16, color: RGBColor = DARK_TEXT):
    """添加项目符号列表文本框。"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {bullet}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = Pt(8)
        p.level = 0
    return tf


def _add_image_safe(slide, image_bytes: bytes, left, top, width, height) -> bool:
    """安全插入图片，失败时返回 False 而非崩溃。"""
    try:
        stream = BytesIO(image_bytes)
        slide.shapes.add_picture(stream, Inches(left), Inches(top), Inches(width), Inches(height))
        return True
    except Exception as e:
        logger.warning(f"插入图片失败: {e}")
        return False


def _add_accent_bar(slide, left, top, width, height, color: RGBColor = ACCENT_COLOR):
    """添加装饰色条（矩形形状，无边框）。"""
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def render_pptx(slides_data: list[dict], title: str, output_path: str):
    """将幻灯片数据列表渲染为 PPTX 文件。

    根据每页的 pageType 分发到对应的渲染函数：
    COVER → _render_cover / CATALOG → _render_catalog
    CONTENT → _render_content / COMPARE → _render_compare / END → _render_end
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for s in slides_data:
        page_type = s.get("pageType", "CONTENT")

        if page_type == "COVER":
            _render_cover(prs, s, title)
        elif page_type == "CATALOG":
            _render_catalog(prs, s)
        elif page_type == "COMPARE":
            _render_compare(prs, s)
        elif page_type == "END":
            _render_end(prs, s)
        else:
            _render_content(prs, s)

    prs.save(output_path)
    logger.info(f"PPTX已生成: {output_path}, 共{len(slides_data)}页")


# ---- 各页面类型渲染函数 ----

def _render_cover(prs, s, title):
    """封面页：深色背景 + 大标题 + 描述 + 装饰线。"""
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, PRIMARY_COLOR)

    _add_accent_bar(slide, 1.5, 3.2, 1.2, 0.06, ACCENT_COLOR)

    _add_textbox(slide, 1.5, 2.2, 10.5, 1.2,
                 s.get("title", title), font_size=44, bold=True, color=WHITE,
                 alignment=PP_ALIGN.LEFT)

    desc = s.get("description", "")
    if desc:
        _add_textbox(slide, 1.5, 3.6, 10.5, 0.8,
                     desc, font_size=20, color=RGBColor(0xBB, 0xBB, 0xCC),
                     alignment=PP_ALIGN.LEFT)

    _add_textbox(slide, 1.5, 5.8, 10.5, 0.5,
                 "AI Generated by Dodo Agent", font_size=12, color=GRAY_TEXT)


def _render_catalog(prs, s):
    """目录页：白色背景 + 左侧装饰线 + 编号章节列表。"""
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)

    _add_accent_bar(slide, 1.0, 1.0, 0.08, 0.6, ACCENT_COLOR)
    _add_textbox(slide, 1.3, 0.9, 8, 0.7,
                 s.get("title", "目录"), font_size=36, bold=True, color=PRIMARY_COLOR)

    items = s.get("items", [])
    if not items and s.get("keyPoints"):
        items = s["keyPoints"]
    if not items and s.get("bullets"):
        items = s["bullets"]

    for i, item in enumerate(items):
        y = 2.2 + i * 1.0
        num = f"0{i + 1}" if i < 9 else str(i + 1)
        _add_textbox(slide, 1.5, y, 0.8, 0.6,
                     num, font_size=32, bold=True, color=ACCENT_COLOR)
        text = item if isinstance(item, str) else item.get("title", str(item))
        _add_textbox(slide, 2.5, y + 0.05, 9, 0.5,
                     text, font_size=22, color=DARK_TEXT)


def _render_content(prs, s):
    """正文内容页：白色背景 + 顶部装饰线 + 标题 + 项目符号 + 可选配图。"""
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)

    _add_accent_bar(slide, 1.0, 0.0, 11.333, 0.04, ACCENT_COLOR)

    _add_textbox(slide, 1.0, 0.5, 11.3, 0.7,
                 s.get("title", ""), font_size=32, bold=True, color=PRIMARY_COLOR)

    subtitle = s.get("subtitle", "")
    if subtitle:
        _add_textbox(slide, 1.0, 1.15, 11.3, 0.5,
                     subtitle, font_size=16, color=GRAY_TEXT)

    has_image = bool(s.get("_image_path") and os.path.exists(s["_image_path"]))
    text_width = 6.5 if has_image else 11.3

    bullets = s.get("bullets", []) or s.get("keyPoints", [])
    if bullets:
        _add_bullet_frame(slide, 1.0, 1.8, text_width, 4.8,
                          bullets, font_size=18, color=DARK_TEXT)

    if has_image:
        with open(s["_image_path"], "rb") as f:
            _add_image_safe(slide, f.read(), 8.0, 1.8, 4.5, 4.5)

    notes = s.get("notes", "")
    if notes:
        _add_textbox(slide, 1.0, 6.8, 11.3, 0.4,
                     notes, font_size=11, color=GRAY_TEXT)


def _render_compare(prs, s):
    """对比页：左右双栏布局，中间 VS 分隔。"""
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)

    _add_accent_bar(slide, 1.0, 0.0, 11.333, 0.04, ACCENT_COLOR)
    _add_textbox(slide, 1.0, 0.5, 11.3, 0.7,
                 s.get("title", ""), font_size=32, bold=True, color=PRIMARY_COLOR)

    left_title = s.get("leftTitle", "方案A")
    right_title = s.get("rightTitle", "方案B")
    left_bullets = s.get("leftBullets", []) or s.get("leftContent", [])
    right_bullets = s.get("rightBullets", []) or s.get("rightContent", [])

    # 左侧栏
    _add_accent_bar(slide, 1.0, 1.6, 5.0, 0.04, ACCENT_COLOR)
    _add_textbox(slide, 1.0, 1.8, 5.0, 0.5,
                 left_title, font_size=24, bold=True, color=ACCENT_COLOR)

    if left_bullets:
        _add_bullet_frame(slide, 1.0, 2.4, 5.0, 4.2,
                          left_bullets, font_size=15)

    # 右侧栏
    _add_accent_bar(slide, 7.0, 1.6, 5.0, 0.04, ACCENT_ORANGE)
    _add_textbox(slide, 7.0, 1.8, 5.0, 0.5,
                 right_title, font_size=24, bold=True, color=ACCENT_ORANGE)

    if right_bullets:
        _add_bullet_frame(slide, 7.0, 2.4, 5.0, 4.2,
                          right_bullets, font_size=15)

    _add_textbox(slide, 6.3, 3.5, 0.8, 0.6,
                 "VS", font_size=28, bold=True, color=GRAY_TEXT,
                 alignment=PP_ALIGN.CENTER)


def _render_end(prs, s):
    """结束页：深色背景 + 居中大标题 + 描述。"""
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, PRIMARY_COLOR)

    _add_textbox(slide, 1.5, 2.8, 10.5, 1.2,
                 s.get("title", "谢谢观看"), font_size=48, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    desc = s.get("description", "")
    if desc:
        _add_textbox(slide, 1.5, 4.2, 10.5, 0.8,
                     desc, font_size=20, color=GRAY_TEXT,
                     alignment=PP_ALIGN.CENTER)


# =============================================================================
# PptBuilderAgent — 状态机驱动的 PPT 自动生成
# =============================================================================

class PptBuilderAgent(BaseAgent):
    """PPT 生成 Agent，五阶段状态机。

    状态流转：INIT → SCHEMA → OUTLINE → CONTENT → RENDER → SUCCESS

    关键特性：
    - 断点续传：每个状态完成后写入 MySQL，中断后可从任意状态恢复
    - Tavily 搜图：为 CONTENT 页自动搜索配图并下载嵌入
    - python-pptx 渲染：5 种页面类型（COVER/CATALOG/CONTENT/COMPARE/END）
    - SSE 实时进度：每页生成进度通过 thinking 事件推送前端
    """

    def __init__(self, conversation_id: str, query: str, file_id: str = ""):
        super().__init__(conversation_id, query, file_id)
        self._inst: AiPptInst | None = None   # PPT 实例（数据库持久化对象）
        self._repo = PptInstRepo()            # 数据访问层

    def _load_or_create_inst(self) -> AiPptInst:
        """加载已有实例（断点续传）或创建新实例。

        - 未完成状态（非 SUCCESS/FAILED）→ 恢复执行
        - 已完成或失败 → 删除旧记录，创建新实例
        """
        existing = self._repo.find_by_conversation_id(self.conversation_id)
        if existing and existing.status not in ("SUCCESS", "FAILED"):
            logger.info(f"恢复PPT任务: {self.conversation_id}, 状态={existing.status}")
            return existing
        if existing and existing.status in ("SUCCESS", "FAILED"):
            self._repo.delete(existing)
            self._repo = PptInstRepo()
        inst = AiPptInst(
            conversation_id=self.conversation_id,
            query=self.query,
            status="INIT",
            create_time=datetime.now(),
            update_time=datetime.now(),
        )
        self._repo.save(inst)
        return inst

    def _save_inst(self, status: str, **fields):
        """更新实例状态和字段，通过 PptInstRepo.update_status() 持久化。

        每个状态处理完成后调用此方法，实现断点续传。
        """
        self._repo.update_status(self._inst, status, **fields)

    async def run(self):
        """状态机主循环：从当前状态开始逐步执行，直至 SUCCESS 或 FAILED。

        流程：
        1. 获取锁 → 加载或创建实例
        2. 从当前状态开始遍历 STATE_ORDER
        3. 每个状态调用对应的 handler 函数
        4. handler 完成后 _save_inst 持久化状态
        5. SUCCESS 时推送 file_ready 事件（含下载链接）
        """
        # ---- 获取锁并注册任务 ----
        ok, error_events = await self._try_start()
        if not ok:
            for evt in error_events:
                yield evt
            return

        llm = build_llm()

        try:
            # ---- 加载或创建实例（断点续传） ----
            self._inst = self._load_or_create_inst()

            start_state = self._inst.status or "INIT"
            if start_state not in STATE_ORDER:
                start_state = "INIT"

            start_idx = STATE_ORDER.index(start_state)
            logger.info(f"PPT Builder 启动: conv={self.conversation_id}, 起始状态={start_state}")

            # ---- 状态机主循环 ----
            for i in range(start_idx, len(STATE_ORDER)):
                state = STATE_ORDER[i]

                # SUCCESS 状态：PPT 已生成，推送下载链接
                if state == "SUCCESS":
                    schema = json.loads(self._inst.ppt_schema or "{}")
                    title = schema.get("title", "output")
                    download_path = f"/agent/pptx/download?conversationId={self.conversation_id}"
                    yield make_event("text",
                                     content=f"\n\nPPT已生成：[点击下载 {title}.pptx]({download_path})")
                    yield make_event("file_ready",
                                     url=self._inst.file_url or "",
                                     fileName=f"{title}.pptx",
                                     message="PPT已生成，点击下载")
                    break

                if self.cancel_event.is_set():
                    raise AgentStopped

                # 根据状态查找对应的 handler
                handler = STATE_HANDLERS.get(state)
                if not handler:
                    self._save_inst("FAILED", error_msg=f"未知状态: {state}")
                    yield make_event("error", message=f"未知状态: {state}")
                    break

                # 执行当前状态的 handler
                try:
                    async for event in handler(self, llm):
                        yield event
                    # handler 内部已通过 _save_inst 持久化，无需额外操作
                except AgentStopped:
                    raise
                except Exception as e:
                    logger.error(f"PPT状态 {state} 异常: {e}")
                    self._save_inst("FAILED", error_msg=str(e))
                    yield make_event("error", message=f"PPT生成失败({state}): {e}")
                    break

        except AgentStopped:
            yield make_sse(json.dumps({"type": "error", "content": "用户已停止"}, ensure_ascii=False))
        except Exception as e:
            logger.error(f"PPT Builder异常: {e}")
            yield make_sse(json.dumps({"type": "error", "content": str(e)}, ensure_ascii=False))
        finally:
            await self._cleanup()


# =============================================================================
# 状态处理器 — 每个状态对应一个 async generator 函数
# =============================================================================

def _thinking(content: str) -> dict:
    """PPT 生成进度通知，末尾加两个换行使前端渲染时与后续内容保持间距。"""
    return make_event("thinking", content=content + "\n\n")


async def _handle_init(agent: PptBuilderAgent, llm):
    """INIT → SCHEMA：分析用户需求，生成 PPT 元信息（标题/页数/风格）。"""
    yield _thinking("正在分析PPT需求...")

    response = await llm.ainvoke([
        ("system", INIT_PROMPT.format(query=agent.query)),
    ])
    text = response.content if hasattr(response, "content") else str(response)
    data = _parse_json(text) or {}

    agent._inst.requirement = data.get("requirement", agent.query)
    agent._inst.template_code = data.get("style", "auto")
    agent._inst.ppt_schema = json.dumps(data, ensure_ascii=False)
    agent._save_inst("SCHEMA")   # 持久化 → 断点续传

    title = data.get("title", "")
    page_count = data.get("pageCount", 8)
    style = data.get("style", "auto")
    yield make_event("thinking",
                     content=f"需求分析完成：主题「{title}」，{page_count}页，风格：{style}")


async def _handle_schema(agent: PptBuilderAgent, llm):
    """SCHEMA → OUTLINE：规划每页的页面类型（COVER/CATALOG/CONTENT/COMPARE/END）。"""
    yield _thinking("正在规划PPT结构...")

    try:
        init_data = json.loads(agent._inst.ppt_schema or "{}")
    except json.JSONDecodeError:
        init_data = {}
    title = init_data.get("title", agent.query[:15])
    page_count = init_data.get("pageCount", 8)
    style = init_data.get("style", "auto")

    response = await llm.ainvoke([
        ("system", SCHEMA_PROMPT.format(
            title=title, requirement=agent._inst.requirement,
            page_count=page_count, style=style)),
    ])
    text = response.content if hasattr(response, "content") else str(response)
    data = _parse_json(text) or {"slides": []}
    slides = data.get("slides", [])

    schema_data = {"title": title, "style": style, "slides": slides}
    agent._inst.ppt_schema = json.dumps(schema_data, ensure_ascii=False)
    agent._save_inst("OUTLINE")

    yield make_event("thinking",
                     content=f"PPT结构规划完成，共 {len(slides)} 页")


async def _handle_outline(agent: PptBuilderAgent, llm):
    """OUTLINE → CONTENT：为每页生成详细大纲要点和配图搜索关键词。"""
    yield _thinking("正在生成详细大纲...")

    schema_json = agent._inst.ppt_schema or "{}"
    try:
        schema_data = json.loads(schema_json)
    except json.JSONDecodeError:
        schema_data = {"title": agent.query, "slides": []}

    response = await llm.ainvoke([
        ("system", OUTLINE_PROMPT.format(
            title=schema_data.get("title", ""),
            schema_json=schema_json)),
    ])
    text = response.content if hasattr(response, "content") else str(response)
    data = _parse_json(text) or {"slides": []}
    outline_slides = data.get("slides", [])

    agent._inst.outline = json.dumps(outline_slides, ensure_ascii=False)
    agent._save_inst("CONTENT")

    yield make_event("thinking",
                     content=f"大纲生成完成，共 {len(outline_slides)} 页")


async def _handle_content(agent: PptBuilderAgent, llm):
    """CONTENT → RENDER：逐页生成详细内容 + Tavily 搜索配图，推送每页进度。

    这是最耗时的阶段，每页需要：
    1. LLM 生成内容（bullets/subtitle/notes）
    2. Tavily 搜索配图（仅 CONTENT 类型页面）
    3. 下载图片并保存到本地临时目录
    """
    try:
        schema_data = json.loads(agent._inst.ppt_schema or "{}")
    except json.JSONDecodeError:
        schema_data = {"title": agent.query, "slides": []}
    try:
        outline_slides = json.loads(agent._inst.outline or "[]")
    except json.JSONDecodeError:
        outline_slides = []

    theme = schema_data.get("title", "")
    schema_slides = schema_data.get("slides", [])
    total = len(schema_slides)

    yield _thinking(f"正在填充幻灯片内容，共 {total} 页...")

    full_slides: list[dict] = []
    for i, schema_slide in enumerate(schema_slides):
        if agent.cancel_event.is_set():
            raise AgentStopped

        slide_index = schema_slide.get("index", i + 1)
        page_type = schema_slide.get("pageType", "CONTENT")
        slide_title = schema_slide.get("title", "")

        # 推送每页生成进度
        yield make_event("thinking",
                         content=f"[{i + 1}/{total}] 生成中：{slide_title}")

        # 匹配大纲中对应页的数据
        outline_match = next(
            (o for o in outline_slides if o.get("index") == slide_index), {})
        key_points = outline_match.get("keyPoints", [])
        image_kw = outline_match.get("imageKeywords", "")

        slide_data = {
            "index": slide_index,
            "pageType": page_type,
            "title": slide_title,
            "description": schema_slide.get("description", ""),
        }

        # CONTENT/COMPARE 页：调用 LLM 生成详细内容
        if page_type in ("CONTENT", "COMPARE"):
            response = await llm.ainvoke([
                ("system", CONTENT_PROMPT.format(
                    theme=theme, page_type=page_type,
                    title=slide_title,
                    key_points=json.dumps(key_points, ensure_ascii=False))),
            ])
            text = response.content if hasattr(response, "content") else str(response)
            content_data = _parse_json(text) or {}

            slide_data["bullets"] = content_data.get("bullets", key_points)
            slide_data["subtitle"] = content_data.get("subtitle", "")
            slide_data["notes"] = content_data.get("notes", "")
            if page_type == "COMPARE":
                slide_data["leftTitle"] = content_data.get("leftTitle", "")
                slide_data["leftBullets"] = content_data.get("leftBullets", [])
                slide_data["rightTitle"] = content_data.get("rightTitle", "")
                slide_data["rightBullets"] = content_data.get("rightBullets", [])
        elif page_type == "CATALOG":
            items = outline_match.get("keyPoints", [])
            slide_data["items"] = [str(item) for item in items] if items else [slide_title]
        elif page_type == "COVER":
            slide_data["title"] = theme
            slide_data["description"] = schema_slide.get("description", "")
        elif page_type == "END":
            slide_data["title"] = "谢谢观看"
            slide_data["description"] = theme

        # 搜索配图（仅 CONTENT 页）
        if image_kw and page_type in ("CONTENT",):
            images = await _search_images(image_kw, max_count=1)
            if images:
                img_bytes = await _download_image(images[0])
                if img_bytes:
                    img_dir = Path(get_settings().upload_dir) / "pptx_images"
                    img_dir.mkdir(parents=True, exist_ok=True)
                    img_name = f"{uuid.uuid4().hex[:12]}.png"
                    img_path = img_dir / img_name
                    img_path.write_bytes(img_bytes)
                    slide_data["_image_path"] = str(img_path)

        full_slides.append(slide_data)

        yield make_event("thinking",
                         content=f"[{i + 1}/{total}] 已完成：{slide_title}")

    # 持久化完整内容数据 → 断点续传
    agent._save_inst("RENDER",
                     ppt_schema=json.dumps({
                         "title": theme,
                         "style": schema_data.get("style", "auto"),
                         "slides": schema_slides,
                     }, ensure_ascii=False),
                     outline=json.dumps(full_slides, ensure_ascii=False))

    yield make_event("thinking",
                     content=f"内容填充完成，共 {len(full_slides)} 页")


async def _handle_render(agent: PptBuilderAgent, llm):
    """RENDER → SUCCESS：调用 python-pptx 渲染 PPTX 文件，上传到 MinIO/本地。"""
    yield _thinking("正在渲染PPT文件...")

    full_slides = json.loads(agent._inst.outline or "[]")
    ppt_schema = json.loads(agent._inst.ppt_schema or "{}")
    title = ppt_schema.get("title", agent.query)

    # 生成 PPTX 文件
    output_dir = Path(get_settings().upload_dir) / "pptx"
    output_dir.mkdir(parents=True, exist_ok=True)
    file_name = f"{uuid.uuid4().hex[:12]}.pptx"
    output_path = output_dir / file_name

    # 同步渲染在线程池中执行，避免阻塞事件循环
    await asyncio.to_thread(render_pptx, full_slides, title, str(output_path))

    # 清理临时图片文件
    for s in full_slides:
        img_path = s.get("_image_path", "")
        if img_path and os.path.exists(img_path):
            try:
                os.remove(img_path)
            except OSError:
                pass

    # 上传到 MinIO（优先）或本地路径
    from minio import Minio
    from minio.error import S3Error

    file_url = ""
    if get_settings().minio_endpoint:
        try:
            client = Minio(
                get_settings().minio_endpoint,
                access_key=get_settings().minio_access_key,
                secret_key=get_settings().minio_secret_key,
                secure=False,
            )
            obj_name = f"pptx/{file_name}"
            client.fput_object(get_settings().minio_bucket, obj_name, str(output_path))
            file_url = f"minio://{get_settings().minio_bucket}/{obj_name}"
        except S3Error as e:
            logger.error(f"MinIO上传PPT失败: {e}")

    if not file_url:
        file_url = f"local://{output_path}"

    agent._save_inst("SUCCESS", file_url=file_url)

    yield _thinking("PPT文件渲染完成")


# 状态 → handler 映射表
STATE_HANDLERS = {
    "INIT": _handle_init,
    "SCHEMA": _handle_schema,
    "OUTLINE": _handle_outline,
    "CONTENT": _handle_content,
    "RENDER": _handle_render,
}
